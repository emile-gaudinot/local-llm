import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
import time
from ollama import chat, ChatResponse
import json
import magic
from PyPDF2 import PdfReader
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation


def get_file_type(file_path):
    if file_path == '':
        return 'no file attached'
    mime = magic.Magic(mime=True)
    file_type = mime.from_file(file_path)
    return file_type

def pdf_content(file_path):
    reader = PdfReader(file_path)
    text = "This is the PDF content parsed twice. Once to extract text, then with OCR:\n\nText content:\n"
    for page in reader.pages:
        text += page.extract_text() + '\n'
    images = convert_from_path(file_path)
    text += "\nOCR content:\n"
    for img in images:
        text += pytesseract.image_to_string(img) + '\n'
    return text

def docx_content(file_path):
    doc = Document(file_path)
    text = "This is the DOCX content:\n\nText content:\n"
    # Extract text
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    # Extract images and perform OCR
    text += "\nImage content:\n"
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img = Image.open(io.BytesIO(rel.target_part.blob))
            img_text = pytesseract.image_to_string(img)
            text += img_text + "\n"
    return text

def pptx_content(file_path):
    prs = Presentation(file_path)
    text = "This is the PPTX content:\n\nText content:\n"
    # Extract text from slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    # Extract images and perform OCR
    text += "\nImage content:\n"
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 is the shape type for pictures
                img = Image.open(io.BytesIO(shape.image.blob))
                img_text = pytesseract.image_to_string(img)
                text += img_text + "\n"
    return text

def update_messages(messages, prompt, file_path):
    file_type = get_file_type(file_path)
    if file_type == 'application/pdf':
        pdf_text = pdf_content(file_path)
        messages += [
            {'role': 'user', 'content': prompt + pdf_text}
        ]
    elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        docx_text = docx_content(file_path)
        messages += [
            {'role': 'user', 'content': prompt + docx_text}
        ]
    elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        pptx_text = pptx_content(file_path)
        messages += [
            {'role': 'user', 'content': prompt + pptx_text}
        ]
    elif file_type in ['image/png', 'image/jpeg', 'image/jpg']:
        # open the image file
        messages += [
            {'role': 'user', 'content': prompt, 'images': [file_path]},
        ]
    else:
        messages += [
            {'role': 'user', 'content': prompt},
        ]
    return messages

def main(model_name, messages_path, file_path):
    elapsed_time = time.time()

    prompt = sys.stdin.read()
    # Load previous messages if exists
    with open(messages_path, 'r', encoding='utf-8') as f:
        messages = json.load(f)
    # Update messages with the prompt and file content
    update_messages(messages, prompt, file_path)

    full_response = ""  
    response: ChatResponse = chat(  
        model=model_name,  
        messages=messages,  
        stream=True
    )
    for chunk in response:
        if chunk.message:  
            response_chunk = chunk.message.content  
            print(response_chunk, end='', flush=True)  
            full_response += response_chunk  
    # Add the exchange to the conversation history 
    messages += [  
        {'role': 'assistant', 'content': full_response},  
    ]  
    # Write messages to temp file
    with open(messages_path, 'w', encoding='utf-8') as f:
        json.dump(messages[-2:], f, ensure_ascii=False, indent=2)

    # Print time
    elapsed_time = time.time() - elapsed_time
    minutes, seconds = divmod(elapsed_time, 60)
    execution_time = f"\n{int(minutes)}min {int(seconds)}s"
    print(execution_time, ' | ', model_name)


if __name__ == '__main__':
    assert len(sys.argv) == 4, f"{sys.argv}"
    model_name, messages_path, file_path = sys.argv[1], sys.argv[2], sys.argv[3]
    main(model_name, messages_path, file_path)